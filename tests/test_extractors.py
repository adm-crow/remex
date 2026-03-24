import io
import json
from pathlib import Path

import pytest

from synapse_core.extractors import (
    extract, extract_metadata, extract_streaming, is_supported, supports_streaming,
)


# --- helpers ---

def write_temp(tmp_path: Path, suffix: str, content: str) -> Path:
    path = tmp_path / f"file{suffix}"
    path.write_text(content, encoding="utf-8")
    return path


# --- txt / md ---

def test_extract_txt(tmp_path):
    path = write_temp(tmp_path, ".txt", "Hello from txt")
    assert extract(path) == "Hello from txt"


def test_extract_md(tmp_path):
    path = write_temp(tmp_path, ".md", "# Title\nSome content")
    assert "Title" in extract(path)


# --- csv ---

def test_extract_csv(tmp_path):
    path = write_temp(tmp_path, ".csv", "name,age\nAlice,30\nBob,25")
    result = extract(path)
    assert "Alice" in result
    assert "Bob" in result


# --- pdf ---

def test_extract_pdf(tmp_path):
    pytest.importorskip("pypdf")
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)

    path = tmp_path / "file.pdf"
    path.write_bytes(buf.getvalue())

    result = extract(path)
    assert isinstance(result, str)


# --- docx ---

def test_extract_docx(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello from docx")
    path = tmp_path / "file.docx"
    doc.save(str(path))

    assert "Hello from docx" in extract(path)


# --- json ---

def test_extract_json(tmp_path):
    data = {"title": "RAG guide", "body": "Retrieval Augmented Generation", "year": 2024}
    path = write_temp(tmp_path, ".json", json.dumps(data))
    result = extract(path)
    assert "RAG guide" in result
    assert "Retrieval Augmented Generation" in result


def test_extract_jsonl(tmp_path):
    lines = [
        json.dumps({"text": "First document"}),
        json.dumps({"text": "Second document"}),
    ]
    path = write_temp(tmp_path, ".jsonl", "\n".join(lines))
    result = extract(path)
    assert "First document" in result
    assert "Second document" in result


# --- html ---

def test_extract_html(tmp_path):
    pytest.importorskip("bs4")
    html = "<html><head><title>Test</title></head><body><p>Hello from HTML</p></body></html>"
    path = write_temp(tmp_path, ".html", html)
    result = extract(path)
    assert "Hello from HTML" in result
    assert "<p>" not in result  # tags stripped


def test_extract_htm_alias(tmp_path):
    pytest.importorskip("bs4")
    html = "<html><body><p>HTM content</p></body></html>"
    path = write_temp(tmp_path, ".htm", html)
    assert "HTM content" in extract(path)


def test_extract_html_strips_scripts(tmp_path):
    pytest.importorskip("bs4")
    html = "<html><body><script>alert('x')</script><p>Real content</p></body></html>"
    path = write_temp(tmp_path, ".html", html)
    result = extract(path)
    assert "Real content" in result
    assert "alert" not in result


# --- pptx ---

def test_extract_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "Hello from PPTX"
    path = tmp_path / "file.pptx"
    prs.save(str(path))

    assert "Hello from PPTX" in extract(path)


# --- xlsx ---

def test_extract_xlsx(tmp_path):
    pytest.importorskip("openpyxl")
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    path = tmp_path / "file.xlsx"
    wb.save(str(path))

    result = extract(path)
    assert "Alice" in result
    assert "Name" in result


# --- epub ---

def test_extract_epub(tmp_path):
    pytest.importorskip("ebooklib")
    pytest.importorskip("bs4")
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_title("Test Book")
    chapter = epub.EpubHtml(title="Chapter 1", file_name="chapter1.xhtml")
    chapter.content = b"<html><body><p>Hello from EPUB</p></body></html>"
    book.add_item(chapter)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]
    path = tmp_path / "file.epub"
    epub.write_epub(str(path), book)

    assert "Hello from EPUB" in extract(path)


# --- odt ---

def test_extract_odt(tmp_path):
    pytest.importorskip("odf")
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    doc = OpenDocumentText()
    p = P()
    p.addText("Hello from ODT")
    doc.text.addElement(p)  # type: ignore[attr-defined]
    path = tmp_path / "file.odt"
    doc.save(str(path))

    assert "Hello from ODT" in extract(path)


# --- unsupported ---

def test_unsupported_extension_raises(tmp_path):
    path = write_temp(tmp_path, ".xyz", "data")
    with pytest.raises(ValueError, match="Unsupported"):
        extract(path)


def test_doc_is_not_supported():
    assert not is_supported(Path("file.doc"))


def test_is_supported():
    assert is_supported(Path("file.txt"))
    assert is_supported(Path("file.PDF"))   # case-insensitive
    assert is_supported(Path("file.json"))
    assert is_supported(Path("file.jsonl"))
    assert is_supported(Path("file.html"))
    assert is_supported(Path("file.htm"))
    assert is_supported(Path("file.pptx"))
    assert is_supported(Path("file.xlsx"))
    assert is_supported(Path("file.epub"))
    assert is_supported(Path("file.odt"))
    assert not is_supported(Path("file.mp3"))


# --- extract_metadata ---

def test_extract_metadata_returns_dict_with_expected_keys(tmp_path):
    path = write_temp(tmp_path, ".txt", "plain text")
    meta = extract_metadata(path)
    assert set(meta.keys()) == {"doc_title", "doc_author", "doc_created"}


def test_extract_metadata_unknown_format_returns_empty_strings(tmp_path):
    path = write_temp(tmp_path, ".txt", "plain text")
    meta = extract_metadata(path)
    assert meta["doc_title"] == ""
    assert meta["doc_author"] == ""
    assert meta["doc_created"] == ""


def test_extract_metadata_pdf(tmp_path):
    pytest.importorskip("pypdf")
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_metadata({"/Title": "My PDF", "/Author": "Test Author"})
    buf = io.BytesIO()
    writer.write(buf)
    path = tmp_path / "file.pdf"
    path.write_bytes(buf.getvalue())

    meta = extract_metadata(path)
    assert meta["doc_title"] == "My PDF"
    assert meta["doc_author"] == "Test Author"


def test_extract_metadata_html(tmp_path):
    pytest.importorskip("bs4")
    html = (
        '<html><head>'
        '<title>My Page</title>'
        '<meta name="author" content="Jane Doe">'
        '</head><body></body></html>'
    )
    path = write_temp(tmp_path, ".html", html)
    meta = extract_metadata(path)
    assert meta["doc_title"] == "My Page"
    assert meta["doc_author"] == "Jane Doe"


def test_extract_metadata_never_raises(tmp_path):
    """extract_metadata must return empty strings even for corrupt files."""
    path = tmp_path / "corrupt.pdf"
    path.write_bytes(b"not a real pdf")
    meta = extract_metadata(path)
    assert isinstance(meta, dict)
    assert all(isinstance(v, str) for v in meta.values())


# =============================================================================
# supports_streaming / extract_streaming
# =============================================================================

def test_supports_streaming_text_formats(tmp_path):
    for suffix in (".txt", ".md", ".csv", ".jsonl"):
        assert supports_streaming(tmp_path / f"file{suffix}") is True


def test_supports_streaming_binary_formats(tmp_path):
    for suffix in (".pdf", ".docx", ".xlsx", ".epub", ".odt", ".json", ".html", ".pptx"):
        assert supports_streaming(tmp_path / f"file{suffix}") is False


def test_extract_streaming_txt_yields_all_content(tmp_path):
    content = "word " * 500  # ~2500 chars
    path = write_temp(tmp_path, ".txt", content)
    blocks = list(extract_streaming(path, block_size=512))
    assert len(blocks) > 1  # must yield multiple blocks for this size
    assert "".join(blocks) == content


def test_extract_streaming_txt_small_file_single_block(tmp_path):
    path = write_temp(tmp_path, ".txt", "short text")
    blocks = list(extract_streaming(path, block_size=65536))
    assert blocks == ["short text"]


def test_extract_streaming_md_yields_all_content(tmp_path):
    content = "# Title\n" + "paragraph " * 200
    path = write_temp(tmp_path, ".md", content)
    rejoined = "".join(extract_streaming(path, block_size=256))
    assert rejoined == content


def test_extract_streaming_csv_yields_all_rows(tmp_path):
    rows = [f"col1_{i}, col2_{i}" for i in range(100)]
    path = write_temp(tmp_path, ".csv", "\n".join(rows))
    blocks = list(extract_streaming(path, block_size=128))
    full = "\n".join(blocks)
    # Every row value must appear somewhere in the output
    assert "col1_0" in full
    assert "col1_99" in full


def test_extract_streaming_jsonl_yields_all_records(tmp_path):
    lines = [f'{{"id": {i}, "text": "record {i}"}}' for i in range(50)]
    path = write_temp(tmp_path, ".jsonl", "\n".join(lines))
    blocks = list(extract_streaming(path, block_size=128))
    full = "\n".join(blocks)
    assert "record 0" in full
    assert "record 49" in full


def test_extract_streaming_jsonl_skips_bad_lines(tmp_path):
    """Invalid JSON lines are skipped; valid lines must still appear in output."""
    content = '{"ok": "yes"}\nnot valid json\n{"also": "fine"}'
    path = write_temp(tmp_path, ".jsonl", content)
    blocks = list(extract_streaming(path))
    full = "".join(blocks)
    assert "yes" in full
    assert "fine" in full
    # The bad line text itself must not leak into the output
    assert "not valid json" not in full


def test_extract_streaming_non_streaming_format_falls_back(tmp_path):
    """Unsupported streaming formats must fall back to a single full-text block."""
    content = '{"key": "value"}'
    path = write_temp(tmp_path, ".json", content)
    blocks = list(extract_streaming(path))
    assert len(blocks) == 1
    assert "value" in blocks[0]


# --- metadata: epub ---

def test_extract_metadata_epub(tmp_path):
    pytest.importorskip("ebooklib")
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_title("My EPUB Title")
    book.add_author("Jane Doe")
    chapter = epub.EpubHtml(title="Ch1", file_name="ch1.xhtml")
    chapter.content = b"<html><body><p>hello</p></body></html>"
    book.add_item(chapter)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", chapter]
    path = tmp_path / "file.epub"
    epub.write_epub(str(path), book)

    meta = extract_metadata(path)
    assert meta["doc_title"] == "My EPUB Title"
    assert meta["doc_author"] == "Jane Doe"


# --- metadata: odt ---

def test_extract_metadata_odt(tmp_path):
    pytest.importorskip("odf")
    from odf.opendocument import OpenDocumentText
    from odf.dc import Title, Creator

    doc = OpenDocumentText()
    doc.meta.addElement(Title(text="My ODT Title"))
    doc.meta.addElement(Creator(text="John Smith"))
    path = tmp_path / "file.odt"
    doc.save(str(path))

    meta = extract_metadata(path)
    assert meta["doc_title"] == "My ODT Title"
    assert meta["doc_author"] == "John Smith"
