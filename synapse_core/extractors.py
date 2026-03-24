import csv
import json
from pathlib import Path
from typing import Iterator

from .logger import logger


def extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("Install pypdf: pip install pypdf")
    reader = PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_csv(path: Path) -> str:
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        return "\n".join(", ".join(row) for row in reader)


def _flatten_json(obj) -> str:
    """Recursively extract all string values from a parsed JSON object."""
    if isinstance(obj, str):
        return obj
    if isinstance(obj, dict):
        return " ".join(_flatten_json(v) for v in obj.values() if v is not None)
    if isinstance(obj, list):
        return " ".join(_flatten_json(item) for item in obj)
    return str(obj) if obj is not None else ""


def extract_json(path: Path) -> str:
    with open(path, encoding="utf-8", errors="ignore") as f:
        data = json.load(f)
    return _flatten_json(data)


def extract_jsonl(path: Path) -> str:
    parts = []
    with open(path, encoding="utf-8", errors="ignore") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                parts.append(_flatten_json(json.loads(line)))
            except json.JSONDecodeError:
                logger.warning("Skipping invalid JSON line in %s", path.name)
    return "\n".join(parts)


def extract_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("Install beautifulsoup4: pip install synapse-core[formats]")
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "head"]):
        tag.decompose()
    return soup.get_text(separator=" ", strip=True)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install synapse-core[formats]")
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():  # type: ignore[attr-defined]
                texts.append(shape.text.strip())  # type: ignore[attr-defined]
    return "\n".join(texts)


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        raise ImportError("Install openpyxl: pip install synapse-core[formats]")
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = ", ".join(str(c) for c in row if c is not None)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


def extract_epub(path: Path) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError(
            "Install ebooklib and beautifulsoup4: pip install synapse-core[formats]"
        )
    book = epub.read_epub(str(path))
    texts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "html.parser")
        text = soup.get_text(separator=" ", strip=True)
        if text:
            texts.append(text)
    return "\n".join(texts)


def extract_odt(path: Path) -> str:
    try:
        from odf.opendocument import load
        from odf.text import P
    except ImportError:
        raise ImportError("Install odfpy: pip install synapse-core[formats]")
    doc = load(str(path))
    parts = []
    for para in doc.getElementsByType(P):
        text = "".join(
            node.data for node in para.childNodes if hasattr(node, "data")
        )
        if text.strip():
            parts.append(text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Streaming extraction (text-based formats only)
# ---------------------------------------------------------------------------

_STREAMING_EXTENSIONS: frozenset[str] = frozenset({".txt", ".md", ".csv", ".jsonl"})


def supports_streaming(path: Path) -> bool:
    """Return True if the file can be extracted page-by-page without full memory load."""
    return path.suffix.lower() in _STREAMING_EXTENSIONS


def extract_streaming(path: Path, block_size: int = 65_536) -> Iterator[str]:
    """Yield text blocks from a file without loading the entire content into memory.

    Supported for ``.txt``, ``.md``, ``.csv``, and ``.jsonl``.
    For any other format this function falls back to :func:`extract` and
    yields the full text as a single block (same semantics, no memory benefit).

    Args:
        path:       File to read.
        block_size: Target characters per block for text-based formats.
                    CSV/JSONL yield batches of rows whose total character count
                    approaches this threshold.
    """
    suffix = path.suffix.lower()

    if suffix in (".txt", ".md"):
        with open(path, encoding="utf-8", errors="ignore") as f:
            while True:
                block = f.read(block_size)
                if not block:
                    break
                yield block

    elif suffix == ".csv":
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            batch: list[str] = []
            size = 0
            for row in reader:
                line = ", ".join(row)
                batch.append(line)
                size += len(line) + 1
                if size >= block_size:
                    yield "\n".join(batch)
                    batch, size = [], 0
            if batch:
                yield "\n".join(batch)

    elif suffix == ".jsonl":
        with open(path, encoding="utf-8", errors="ignore") as f:
            batch_j: list[str] = []
            size_j = 0
            for raw_line in f:
                raw_line = raw_line.strip()
                if not raw_line:
                    continue
                try:
                    line = _flatten_json(json.loads(raw_line))
                except json.JSONDecodeError:
                    logger.warning("Skipping invalid JSON line in %s", path.name)
                    continue
                batch_j.append(line)
                size_j += len(line) + 1
                if size_j >= block_size:
                    yield "\n".join(batch_j)
                    batch_j, size_j = [], 0
            if batch_j:
                yield "\n".join(batch_j)

    else:
        # Binary / non-streaming format — load once, yield once
        yield extract(path)


EXTRACTORS = {
    ".txt":  extract_txt,
    ".md":   extract_txt,
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".csv":  extract_csv,
    ".json": extract_json,
    ".jsonl":extract_jsonl,
    ".html": extract_html,
    ".htm":  extract_html,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".epub": extract_epub,
    ".odt":  extract_odt,
}


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    extractor = EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix}")
    return extractor(path)


SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(EXTRACTORS.keys())
"""Frozenset of file extensions (with leading dot) that synapse-core can ingest."""


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in EXTRACTORS


def extract_metadata(path: Path) -> dict[str, str]:
    """Extract document-level metadata (title, author, created) where available.

    Returns a dict with keys ``doc_title``, ``doc_author``, ``doc_created``.
    Missing or unreadable fields are returned as empty strings.
    Guaranteed never to raise.
    """
    meta: dict[str, str] = {"doc_title": "", "doc_author": "", "doc_created": ""}
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
            info = PdfReader(str(path)).metadata
            if info:
                meta["doc_title"] = info.title or ""
                meta["doc_author"] = info.author or ""
                if info.creation_date:
                    meta["doc_created"] = info.creation_date.isoformat()
        except Exception as e:
            logger.debug("Could not read PDF metadata from %s: %s", path.name, e)

    elif suffix == ".docx":
        try:
            from docx import Document
            props = Document(str(path)).core_properties
            meta["doc_title"] = props.title or ""
            meta["doc_author"] = props.author or ""
            if props.created:
                meta["doc_created"] = props.created.isoformat()
        except Exception as e:
            logger.debug("Could not read DOCX metadata from %s: %s", path.name, e)

    elif suffix in (".html", ".htm"):
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(
                path.read_text(encoding="utf-8", errors="ignore"), "html.parser"
            )
            if soup.title and soup.title.string:
                meta["doc_title"] = soup.title.string.strip()
            author_tag = soup.find("meta", attrs={"name": "author"})
            if author_tag:
                meta["doc_author"] = author_tag.get("content", "")  # type: ignore[arg-type]
        except Exception as e:
            logger.debug("Could not read HTML metadata from %s: %s", path.name, e)

    elif suffix == ".pptx":
        try:
            from pptx import Presentation
            props = Presentation(str(path)).core_properties
            meta["doc_title"] = props.title or ""
            meta["doc_author"] = props.author or ""
            if props.created:
                meta["doc_created"] = props.created.isoformat()
        except Exception as e:
            logger.debug("Could not read PPTX metadata from %s: %s", path.name, e)

    elif suffix == ".epub":
        try:
            from ebooklib import epub
            book = epub.read_epub(str(path))
            titles = book.get_metadata("DC", "title")
            if titles:
                meta["doc_title"] = titles[0][0] or ""
            creators = book.get_metadata("DC", "creator")
            if creators:
                meta["doc_author"] = creators[0][0] or ""
            dates = book.get_metadata("DC", "date")
            if dates:
                meta["doc_created"] = dates[0][0] or ""
        except Exception as e:
            logger.debug("Could not read EPUB metadata from %s: %s", path.name, e)

    elif suffix == ".odt":
        try:
            from odf.opendocument import load
            from odf.dc import Title, Creator, Date
            doc = load(str(path))
            titles = doc.meta.getElementsByType(Title)
            if titles:
                meta["doc_title"] = str(titles[0])
            creators = doc.meta.getElementsByType(Creator)
            if creators:
                meta["doc_author"] = str(creators[0])
            dates = doc.meta.getElementsByType(Date)
            if dates:
                meta["doc_created"] = str(dates[0])
        except Exception as e:
            logger.debug("Could not read ODT metadata from %s: %s", path.name, e)

    return meta
